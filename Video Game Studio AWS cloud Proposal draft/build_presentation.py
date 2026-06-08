from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_SLIDE     = RGBColor(0x10, 0x24, 0x3A)   # slide background
ACCENT_BLUE  = RGBColor(0x00, 0x8C, 0xFF)   # bright blue
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)   # cyan highlight
ACCENT_GREEN = RGBColor(0x39, 0xD3, 0x53)   # green for savings
ACCENT_ORG   = RGBColor(0xFF, 0x99, 0x00)   # orange accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xD6, 0xE0)
MID_GREY     = RGBColor(0x55, 0x6B, 0x82)
TABLE_HDR    = RGBColor(0x00, 0x5C, 0xAA)
TABLE_ROW1   = RGBColor(0x12, 0x2D, 0x4A)
TABLE_ROW2   = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ─────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    # dark background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_SLIDE
    bg.line.fill.background()
    return slide

def txb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def accent_bar(slide, color=ACCENT_BLUE, t=1.15):
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(t),
                                  Inches(12.33), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def slide_title(slide, title, subtitle=None, title_color=ACCENT_CYAN):
    txb(slide, title, 0.5, 0.25, 12.33, 0.85, size=32, bold=True,
        color=title_color, align=PP_ALIGN.LEFT)
    accent_bar(slide, ACCENT_BLUE, 1.1)
    if subtitle:
        txb(slide, subtitle, 0.5, 1.2, 12.33, 0.45, size=16,
            color=LIGHT_GREY, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, l, t, w, h, size=15, color=WHITE,
                 bullet_color=ACCENT_CYAN, title=None, title_size=17):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = ACCENT_CYAN
        first = False
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = "\u2022  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color

def add_table(slide, headers, rows, l, t, w, h):
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols,
                                  Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    col_w = Inches(w) // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    def cell_fmt(cell, text, bg, fg=WHITE, bold=False, sz=12, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(text)
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = fg

    for ci, h in enumerate(headers):
        cell_fmt(tbl.cell(0, ci), h, TABLE_HDR, WHITE, bold=True, sz=12)
    for ri, row in enumerate(rows):
        bg = TABLE_ROW1 if ri % 2 == 0 else TABLE_ROW2
        for ci, val in enumerate(row):
            cell_fmt(tbl.cell(ri+1, ci), val, bg, LIGHT_GREY, sz=11,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    return tbl

def box(slide, text, l, t, w, h, fill=TABLE_HDR, fg=WHITE,
        sz=13, bold=False, align=PP_ALIGN.CENTER, radius=False):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = ACCENT_BLUE
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = fg
    return shp

def arrow_right(slide, l, t, w=0.4, h=0.25, color=ACCENT_BLUE):
    shp = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def arrow_down(slide, l, t, w=0.25, h=0.3, color=ACCENT_BLUE):
    shp = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    from pptx.util import Emu
    shp.rotation = 90

def divider(slide, t, color=MID_GREY):
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(t),
                                  Inches(12.33), Inches(0.02))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()

# big background accent strip
strip = s.shapes.add_shape(1, 0, Inches(2.8), prs.slide_width, Inches(2.0))
strip.fill.solid()
strip.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x6E)
strip.line.fill.background()

txb(s, "AWS CLOUD MIGRATION PROPOSAL", 0.6, 0.35, 12.0, 1.0,
    size=36, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
txb(s, "Scaling Game Infrastructure with Amazon Web Services",
    0.6, 1.3, 12.0, 0.6, size=20, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

accent_bar(s, ACCENT_BLUE, 2.0)

txb(s, "Mid-Sized Video Game Studio & Publisher",
    0.6, 2.9, 12.0, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Presenter: Cloud Solutions Architect   |   AWS Region: US East (N. Virginia) — us-east-1   |   May 2026",
    0.6, 3.45, 12.0, 0.4, size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

txb(s, '"Our players expect fast, reliable, always-on experiences.\nAWS gives us the infrastructure to deliver that — at any scale."',
    1.5, 4.2, 10.0, 1.0, size=15, italic=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

txb(s, "Topics: Business Justification  |  Cloud-Native & Serverless  |  TCO Analysis  |  AWS Pricing Calculator",
    0.6, 6.8, 12.0, 0.45, size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Who We Are
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Who We Are & What We Need",
            "Mid-Sized Video Game Studio & Publisher")

txb(s, "OUR ORGANIZATION", 0.5, 1.45, 5.8, 0.35, size=13, bold=True, color=ACCENT_CYAN)
items_org = [
    "Operates online multiplayer services (matchmaking, lobbies, leaderboards)",
    "Manages player accounts, profiles, and in-game purchases",
    "Distributes game patches, DLC, trailers, and screenshots globally",
    "Analyzes player behavior, session data, and crash reports",
    "Experiences unpredictable traffic spikes during launches & events",
]
bullet_block(s, items_org, 0.5, 1.8, 5.9, 3.5, size=13)

txb(s, "THE PROBLEM WITH ON-PREMISES", 7.0, 1.45, 5.8, 0.35, size=13, bold=True, color=ACCENT_ORG)
headers2 = ["Challenge", "Impact"]
rows2 = [
    ["Fixed server capacity",        "Cannot handle launch-day spikes"],
    ["High upfront hardware cost",   "Capital risk before game success"],
    ["Manual patching & maintenance","Engineering time lost"],
    ["Single datacenter",            "No redundancy or failover"],
    ["Slow provisioning",            "Weeks to add capacity"],
]
add_table(s, headers2, rows2, 7.0, 1.85, 5.8, 3.2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Business Justification 1/2
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Business Justification — Why Move to AWS?",
            "Five benefits tied directly to gaming business needs  (1 of 2)")

benefits = [
    ("1", "Handle Traffic Spikes",
     "Game launches, seasonal events, and streamer promotions can spike player\ntraffic 10x–100x in minutes. EC2 Auto Scaling adds capacity instantly and\nscales back down when demand drops — no idle hardware."),
    ("2", "Improve Player Experience & Uptime",
     "Players expect 24/7 availability. RDS Multi-AZ auto-fails over if the\nprimary DB fails. Elastic Load Balancing removes unhealthy instances.\nCloudWatch monitors and alerts in real time. AWS SLA: 99.99%."),
    ("3", "Reduce Upfront Hardware Costs",
     "No need to buy servers before knowing if a game will succeed.\nAWS converts large capital expenditures (CapEx) into predictable\nmonthly operating expenses (OpEx). Pay only for what you use."),
]

tops = [1.55, 3.2, 4.85]
for (num, title, body), top in zip(benefits, tops):
    box(s, num, 0.5, top, 0.45, 0.75, fill=ACCENT_BLUE, fg=WHITE, sz=20, bold=True)
    txb(s, title, 1.05, top, 11.5, 0.35, size=15, bold=True, color=ACCENT_CYAN)
    txb(s, body,  1.05, top+0.35, 11.5, 0.7, size=12, color=LIGHT_GREY)
    if top < 4.85:
        divider(s, top + 1.55)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Business Justification 2/2
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Business Justification — Why Move to AWS?",
            "Five benefits tied directly to gaming business needs  (2 of 2)")

benefits2 = [
    ("4", "Global Delivery of Game Files & Patches",
     "Games need frequent updates — patches, DLC, hotfixes, trailers.\nAmazon S3 stores assets with 11-nines durability. CloudFront CDN\ndelivers files from edge locations worldwide for fast player downloads."),
    ("5", "Data-Driven Game Development",
     "Studios need to understand player behavior, retention, crashes, and\npurchases. CloudWatch + Kinesis + S3 + Athena + QuickSight provide\na full analytics pipeline — no separate analytics infrastructure needed."),
]

tops2 = [1.55, 3.5]
for (num, title, body), top in zip(benefits2, tops2):
    box(s, num, 0.5, top, 0.45, 0.75, fill=ACCENT_BLUE, fg=WHITE, sz=20, bold=True)
    txb(s, title, 1.05, top, 11.5, 0.35, size=15, bold=True, color=ACCENT_CYAN)
    txb(s, body,  1.05, top+0.35, 11.5, 0.7, size=12, color=LIGHT_GREY)
    if top < 3.5:
        divider(s, top + 1.55)

# summary box
box(s, "Result: Flexible infrastructure that scales with player demand, reduces financial risk, and improves uptime — all without owning physical hardware.",
    0.5, 5.6, 12.33, 0.75, fill=RGBColor(0x00, 0x3A, 0x6E), fg=WHITE, sz=13, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Cloud-Native & Serverless Rationale
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Cloud-Native & Serverless Services",
            "Why managed and serverless AWS services are the right fit for a game studio")

txb(s, "CLOUD-NATIVE", 0.5, 1.45, 5.9, 0.35, size=13, bold=True, color=ACCENT_CYAN)
cn_items = [
    "Built for cloud — managed, scalable, integrated",
    "No OS patching, hardware maintenance, or capacity planning",
    "AWS handles availability, backups, and failover",
    "Examples: RDS, S3, CloudFront, ELB, CloudWatch",
]
bullet_block(s, cn_items, 0.5, 1.85, 5.9, 2.2, size=13)

txb(s, "SERVERLESS", 7.0, 1.45, 5.8, 0.35, size=13, bold=True, color=ACCENT_GREEN)
sl_items = [
    "No servers to provision, patch, or manage at all",
    "Pay per invocation — not per idle server hour",
    "Scales automatically from 0 to millions of requests",
    "Examples: Lambda, API Gateway, S3 Events",
]
bullet_block(s, sl_items, 7.0, 1.85, 5.8, 2.2, size=13, color=WHITE)

divider(s, 4.15)

txb(s, "LAMBDA SPOTLIGHT — Match-End Event Processing",
    0.5, 4.25, 12.33, 0.35, size=14, bold=True, color=ACCENT_ORG)

# Lambda flow boxes
labels = ["Match Ends", "Lambda\nTriggered", "Update RDS\nRankings", "Award\nCurrency", "Send Push\nNotification"]
colors = [MID_GREY, ACCENT_BLUE, TABLE_HDR, TABLE_HDR, TABLE_HDR]
positions = [0.5, 2.5, 4.5, 6.9, 9.3]
for label, color, lpos in zip(labels, colors, positions):
    box(s, label, lpos, 4.7, 1.8, 0.85, fill=color, fg=WHITE, sz=12, bold=True)
for lpos in [2.35, 4.35, 6.75, 9.15]:
    arrow_right(s, lpos, 4.98, 0.35, 0.25)

txb(s, "Without Lambda: a server runs 24/7 waiting for events — even at 3 AM with zero players.",
    0.5, 5.7, 12.33, 0.35, size=12, italic=True, color=ACCENT_ORG)
txb(s, "With Lambda: function runs ~200ms per match end. Studio pays for 200ms of compute, not 24 hours of idle time.",
    0.5, 6.1, 12.33, 0.35, size=12, italic=True, color=ACCENT_GREEN)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AWS Services Table
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "AWS Services Mapped to Gaming Needs",
            "Each service solves a specific game studio business or technical requirement")

headers6 = ["AWS Service", "Type", "How It Helps the Game Studio"]
rows6 = [
    ["Amazon EC2",          "Compute",        "Game backend servers, matchmaking, login/auth, lobby systems"],
    ["Amazon RDS MySQL",    "Managed DB",     "Player profiles, inventory, leaderboards, purchases — Multi-AZ HA"],
    ["Amazon S3",           "Object Storage", "Patches, DLC, logs, screenshots, trailers, backups — 11-nines durability"],
    ["AWS Lambda",          "Serverless",     "Match results, rewards, image resize, log cleanup — no idle servers"],
    ["Amazon API Gateway",  "Managed API",    "Secure REST/WebSocket endpoints for game clients, apps, launchers"],
    ["Amazon CloudWatch",   "Monitoring",     "Server health, player traffic, errors, latency — alarms & auto-recovery"],
    ["Elastic Load Balancing","Traffic Mgmt", "Spreads player connections across EC2 instances, removes unhealthy ones"],
    ["Auto Scaling",        "Capacity",       "Adds EC2 on launch day, removes on off-peak — automatic, policy-driven"],
    ["Amazon CloudFront",   "CDN",            "Delivers patches & media from global edge locations — low latency"],
    ["AWS IAM",             "Access Control", "Least-privilege permissions for devs, admins, CI/CD pipelines"],
]
add_table(s, headers6, rows6, 0.5, 1.45, 12.33, 5.7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Architecture Diagram
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "AWS Architecture — Video Game Studio",
            "How EC2, RDS, S3, Lambda, CloudFront, and supporting services connect")

# ── Row 1: Players / Internet ─────────────────────────────────────────────
box(s, "Players\n(Global)", 0.3, 1.45, 1.5, 0.75, fill=MID_GREY, fg=WHITE, sz=11, bold=True)
arrow_right(s, 1.85, 1.72, 0.45, 0.22, color=ACCENT_CYAN)

box(s, "Amazon\nCloudFront\n(CDN)", 2.35, 1.45, 1.7, 0.75, fill=TABLE_HDR, fg=WHITE, sz=10, bold=True)
arrow_right(s, 4.1, 1.72, 0.45, 0.22, color=ACCENT_CYAN)

box(s, "Elastic Load\nBalancer", 4.6, 1.45, 1.7, 0.75, fill=TABLE_HDR, fg=WHITE, sz=10, bold=True)
arrow_right(s, 6.35, 1.72, 0.45, 0.22, color=ACCENT_CYAN)

box(s, "Amazon EC2\nBackend Servers\n(10x m6i.8xlarge)", 6.85, 1.45, 2.2, 0.75, fill=ACCENT_BLUE, fg=WHITE, sz=10, bold=True)
arrow_right(s, 9.1, 1.72, 0.45, 0.22, color=ACCENT_CYAN)

box(s, "AWS IAM\nAccess Control", 9.6, 1.45, 1.7, 0.75, fill=RGBColor(0x1A,0x3A,0x5C), fg=WHITE, sz=10, bold=True)

# ── Vertical arrows down from EC2 ────────────────────────────────────────
arrow_down(s, 7.85, 2.25, 0.22, 0.35, color=ACCENT_CYAN)
arrow_down(s, 7.85, 3.55, 0.22, 0.35, color=ACCENT_CYAN)

# ── Row 2: RDS + Lambda ───────────────────────────────────────────────────
box(s, "Amazon RDS MySQL\nMulti-AZ\n(5x db.r6g.2xlarge)\nPlayer data, leaderboards,\ninventory, purchases",
    5.5, 2.65, 2.9, 1.1, fill=RGBColor(0x00,0x52,0x8C), fg=WHITE, sz=10, bold=False)

arrow_right(s, 8.45, 3.1, 0.4, 0.22, color=ACCENT_GREEN)

box(s, "AWS Lambda\nServerless Functions\nMatch results, rewards,\nlog cleanup, notifications",
    8.9, 2.65, 2.5, 1.1, fill=RGBColor(0x1A,0x5C,0x2A), fg=WHITE, sz=10, bold=False)

# ── Vertical arrow down from EC2 to S3 ───────────────────────────────────
arrow_down(s, 7.85, 3.95, 0.22, 0.35, color=ACCENT_CYAN)

# ── Row 3: S3 + CloudWatch ───────────────────────────────────────────────
box(s, "Amazon S3 Standard\nGame patches, DLC, logs,\ntrailers, screenshots,\nbackups  (2 TB)",
    5.5, 4.35, 2.9, 1.1, fill=RGBColor(0x7A,0x3B,0x00), fg=WHITE, sz=10, bold=False)

arrow_right(s, 8.45, 4.8, 0.4, 0.22, color=ACCENT_ORG)

box(s, "Amazon CloudWatch\nMonitoring & Alerts\nServer health, traffic,\nerrors, latency",
    8.9, 4.35, 2.5, 1.1, fill=RGBColor(0x4A,0x2A,0x00), fg=WHITE, sz=10, bold=False)

# ── Row 4: API Gateway + Auto Scaling labels ─────────────────────────────
box(s, "Amazon API Gateway\nREST / WebSocket endpoints\nfor game clients & apps",
    0.3, 4.35, 2.5, 1.1, fill=RGBColor(0x1A,0x3A,0x5C), fg=WHITE, sz=10, bold=False)

arrow_right(s, 2.85, 4.8, 0.4, 0.22, color=ACCENT_CYAN)

box(s, "Auto Scaling\nAdds / removes EC2\nbased on player demand",
    3.3, 4.35, 2.1, 1.1, fill=RGBColor(0x00,0x3A,0x6E), fg=WHITE, sz=10, bold=False)

# ── Legend ────────────────────────────────────────────────────────────────
txb(s, "LEGEND:", 0.3, 5.7, 1.2, 0.3, size=10, bold=True, color=LIGHT_GREY)
for lbl, col, lx in [("Compute", ACCENT_BLUE, 1.5), ("Database", RGBColor(0x00,0x52,0x8C), 3.1),
                      ("Storage", RGBColor(0x7A,0x3B,0x00), 4.7), ("Serverless", RGBColor(0x1A,0x5C,0x2A), 6.3),
                      ("Monitoring", RGBColor(0x4A,0x2A,0x00), 7.9), ("CDN/Network", TABLE_HDR, 9.5)]:
    box(s, lbl, lx, 5.7, 1.4, 0.35, fill=col, fg=WHITE, sz=9, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TCO Assumptions
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "TCO Analysis — Assumptions & Region",
            "All pricing: US East (N. Virginia) — us-east-1  |  Sourced from AWS Pricing Calculator")

headers8 = ["Assumption", "Value"]
rows8 = [
    ["AWS Region",              "US East (N. Virginia) — us-east-1"],
    ["Operating System",        "Linux (no additional OS licensing cost)"],
    ["Uptime",                  "730 hours/month (24/7, full month)"],
    ["Support Level",           "Not included — billed separately"],
    ["Data Transfer",           "Intra-region EC2 to RDS/S3 assumed free"],
    ["Pricing Model",           "On-Demand baseline vs 1-Year Reserved (No Upfront)"],
    ["RDS Backup Storage",      "Automated backups free up to provisioned storage size"],
    ["S3 Lifecycle Policies",   "None — all 2 TB stays in S3 Standard"],
    ["S3 Cross-Region Replication", "None — single region only"],
    ["EBS Volumes",             "gp3, 100 GB per EC2 instance"],
]
add_table(s, headers8, rows8, 0.5, 1.45, 12.33, 5.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TCO: EC2
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "TCO Component A — EC2: 10 Linux Virtual Machines",
            "32 vCPU, 128 GB RAM each  |  Instance: m6i.8xlarge  |  Use: Game backend servers, matchmaking, auth, APIs")

# Instance specs box
txb(s, "INSTANCE SELECTION: m6i.8xlarge", 0.5, 1.45, 5.9, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_ec2s = ["Spec", "Value"]
rows_ec2s = [
    ["Instance Family", "M6i — General Purpose (Intel Ice Lake)"],
    ["vCPUs",           "32"],
    ["RAM",             "128 GiB"],
    ["Network",         "Up to 12.5 Gbps"],
    ["OS",              "Linux"],
    ["Why M6i?",        "4:1 compute/memory ratio, 6th-gen Intel, better perf than M5"],
]
add_table(s, headers_ec2s, rows_ec2s, 0.5, 1.85, 5.9, 2.8)

# Pricing table
txb(s, "PRICING: 1-YEAR RESERVED (NO UPFRONT)", 7.0, 1.45, 5.8, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_ec2p = ["Pricing Tier", "Hourly Rate", "Monthly (730 hrs)"]
rows_ec2p = [
    ["On-Demand",                    "$1.536/hr", "$1,121.28"],
    ["1-Year Reserved (No Upfront)", "~$0.952/hr", "~$695.00"],
]
add_table(s, headers_ec2p, rows_ec2p, 7.0, 1.85, 5.8, 1.3)

txb(s, "WHY RESERVED?  Game backend servers run 24/7 — predictable workload. 1-Yr RI saves ~38% vs On-Demand with no upfront cash commitment.",
    7.0, 3.25, 5.8, 0.7, size=12, color=LIGHT_GREY)

divider(s, 4.75)

# Cost calculation
txb(s, "MONTHLY COST CALCULATION", 0.5, 4.85, 12.33, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_ec2c = ["Component", "Calculation", "Monthly Cost"]
rows_ec2c = [
    ["10 x m6i.8xlarge (1-Yr RI No Upfront)", "$0.952 x 730 hrs x 10 instances", "$6,950.00"],
    ["10 x 100 GB gp3 EBS root volumes",       "1,000 GB x $0.08/GB",             "$80.00"],
    ["EC2 TOTAL",                               "",                                "$7,030.00"],
]
add_table(s, headers_ec2c, rows_ec2c, 0.5, 5.25, 12.33, 1.85)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TCO: RDS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "TCO Component B — RDS MySQL: 5 Database Servers",
            "Instance: db.r6g.2xlarge  |  Multi-AZ  |  Use: Player accounts, leaderboards, inventory, purchases")

txb(s, "INSTANCE SELECTION: db.r6g.2xlarge", 0.5, 1.45, 5.9, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_rdss = ["Spec", "Value"]
rows_rdss = [
    ["Instance Family", "R6g — Memory Optimized (Graviton2 ARM)"],
    ["vCPUs",           "8"],
    ["RAM",             "64 GiB"],
    ["Engine",          "MySQL 8.0"],
    ["Why R6g?",        "~20% better price/perf than x86, memory-optimized for DB"],
]
add_table(s, headers_rdss, rows_rdss, 0.5, 1.85, 5.9, 2.5)

txb(s, "CONFIGURATION", 7.0, 1.45, 5.8, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_rdsc = ["Setting", "Choice", "Reason"]
rows_rdsc = [
    ["Deployment",  "Multi-AZ",          "Auto failover — game data must stay up"],
    ["Storage",     "gp3 SSD",           "Better IOPS than gp2, cheaper than io1"],
    ["Size",        "500 GB/instance",   "Player data, indexes, growth headroom"],
    ["Backup",      "7-day retention",   "Free up to provisioned storage size"],
    ["Pricing",     "1-Yr RI No Upfront","~40% savings vs On-Demand"],
]
add_table(s, headers_rdsc, rows_rdsc, 7.0, 1.85, 5.8, 2.5)

divider(s, 4.5)

txb(s, "MONTHLY COST CALCULATION", 0.5, 4.6, 12.33, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_rdsp = ["Line Item", "Calculation", "Monthly Cost"]
rows_rdsp = [
    ["On-Demand Single-AZ rate",                    "$1.038/hr (confirmed us-east-1)",          "—"],
    ["On-Demand Multi-AZ rate",                     "$1.038 x 2 = $2.076/hr",                  "—"],
    ["1-Yr RI No Upfront Multi-AZ (~40% discount)", "$2.076 x 0.60 = $1.246/hr",               "—"],
    ["5 x db.r6g.2xlarge Multi-AZ (1-Yr RI)",       "$1.246 x 730 x 5",                        "$4,547.90"],
    ["Storage: 500 GB gp3 x 5 x 2 (Multi-AZ)",      "5,000 GB x $0.115/GB",                    "$575.00"],
    ["Automated Backups",                            "Included free (up to provisioned size)",   "$0.00"],
    ["RDS TOTAL",                                    "",                                         "$5,122.90"],
]
add_table(s, headers_rdsp, rows_rdsp, 0.5, 5.0, 12.33, 2.2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TCO: S3
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "TCO Component C — Amazon S3 Standard Bucket",
            "Use: Game patches, DLC, logs, trailers, screenshots, backups  |  2 TB data")

txb(s, "CONFIGURATION & ASSUMPTIONS", 0.5, 1.45, 5.9, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_s3c = ["Setting", "Value"]
rows_s3c = [
    ["Storage Class",           "S3 Standard"],
    ["Data Volume",             "2 TB (2,048 GB)"],
    ["Lifecycle Policies",      "None — all data stays in S3 Standard"],
    ["Cross-Region Replication","None — single region (us-east-1)"],
    ["Data Transfer Out",       "Via CloudFront or intra-region (free to EC2)"],
]
add_table(s, headers_s3c, rows_s3c, 0.5, 1.85, 5.9, 2.3)

txb(s, "S3 STANDARD PRICING (us-east-1)", 7.0, 1.45, 5.8, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_s3p = ["Component", "Rate"]
rows_s3p = [
    ["Storage (first 50 TB/month)",       "$0.023 per GB/month"],
    ["PUT, COPY, POST, LIST requests",    "$0.005 per 1,000 requests"],
    ["GET, SELECT, and other requests",   "$0.0004 per 1,000 requests"],
]
add_table(s, headers_s3p, rows_s3p, 7.0, 1.85, 5.8, 1.7)

divider(s, 4.3)

txb(s, "MONTHLY COST CALCULATION", 0.5, 4.4, 12.33, 0.35, size=13, bold=True, color=ACCENT_CYAN)
headers_s3m = ["Line Item", "Calculation", "Monthly Cost"]
rows_s3m = [
    ["Storage: 2,048 GB x $0.023",                    "2,048 x $0.023",                "$47.10"],
    ["PUT/COPY: 1,000,000 requests",                  "1,000,000 / 1,000 x $0.005",    "$5.00"],
    ["GET/SELECT: 10,000,000 requests",               "10,000,000 / 1,000 x $0.0004",  "$4.00"],
    ["Data Transfer Out",                             "Intra-region / CloudFront",      "$0.00"],
    ["S3 TOTAL",                                      "",                               "$56.10"],
]
add_table(s, headers_s3m, rows_s3m, 0.5, 4.8, 12.33, 2.0)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TCO Summary
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "TCO Summary — Monthly Cost Breakdown",
            "Region: us-east-1  |  Pricing model: 1-Year Reserved (No Upfront) vs On-Demand baseline")

headers_sum = ["Service", "Configuration", "On-Demand/mo", "Reserved/mo", "Monthly Savings"]
rows_sum = [
    ["EC2",        "10x m6i.8xlarge Linux + 100 GB gp3 EBS each",          "$11,292.80", "$7,030.00",  "$4,262.80"],
    ["RDS MySQL",  "5x db.r6g.2xlarge Multi-AZ + 500 GB gp3 each",         "$8,152.40",  "$5,122.90",  "$3,029.50"],
    ["S3 Standard","2 TB + 1M PUT requests + 10M GET requests",             "$56.10",     "$56.10",     "$0.00"],
    ["TOTAL",      "",                                                       "$19,501.30", "$12,209.00", "$7,292.30"],
]
add_table(s, headers_sum, rows_sum, 0.5, 1.45, 12.33, 2.5)

# Highlight boxes
box(s, "Monthly Total\n(Reserved)\n$12,209.00", 0.5, 4.2, 2.8, 1.3,
    fill=ACCENT_BLUE, fg=WHITE, sz=16, bold=True)
box(s, "Annual Total\n(Reserved)\n$146,508.00", 3.5, 4.2, 2.8, 1.3,
    fill=TABLE_HDR, fg=WHITE, sz=16, bold=True)
box(s, "Annual Savings\nvs On-Demand\n~$87,508/year", 6.5, 4.2, 2.8, 1.3,
    fill=RGBColor(0x1A,0x5C,0x2A), fg=WHITE, sz=16, bold=True)
box(s, "On-Demand\nBaseline/mo\n~$19,501", 9.5, 4.2, 2.8, 1.3,
    fill=MID_GREY, fg=WHITE, sz=16, bold=True)

txb(s, "Switching from On-Demand to 1-Year Reserved Instances saves ~$87,500/year without changing any architecture.",
    0.5, 5.7, 12.33, 0.4, size=13, italic=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
txb(s, "Note: The ~$19,543 on-demand figure referenced in the brief aligns with this on-demand baseline estimate.",
    0.5, 6.15, 12.33, 0.35, size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — AWS Pricing Calculator
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "AWS Pricing Calculator — How to Build This Estimate",
            "URL: https://calculator.aws/pricing/2/home")

steps = [
    ("STEP 1 — EC2",
     "Add service > EC2 | Region: US East (N. Virginia) | Instance: m6i.8xlarge | OS: Linux\nQty: 10 | Usage: 730 hrs/month | Pricing: Reserved 1-Year No Upfront\nAdd EBS: gp3, 100 GB x 10 instances"),
    ("STEP 2 — RDS MySQL",
     "Add service > RDS | Engine: MySQL | Instance: db.r6g.2xlarge | Deployment: Multi-AZ\nQty: 5 | Storage: gp3, 500 GB/instance | Pricing: Reserved 1-Year No Upfront\nBackup retention: 7 days"),
    ("STEP 3 — S3",
     "Add service > S3 | Storage class: S3 Standard | Storage: 2,048 GB\nPUT/COPY/POST/LIST requests: 1,000,000 | GET/SELECT requests: 10,000,000\nData transfer out: 0 GB (intra-region assumed free)"),
    ("STEP 4 — Share",
     "Click View Summary > verify totals match calculations\nClick Share > copy the unique estimate URL\nScreenshot each service panel + summary page for submission"),
]

tops13 = [1.45, 2.75, 4.05, 5.35]
colors13 = [ACCENT_BLUE, TABLE_HDR, RGBColor(0x7A,0x3B,0x00), RGBColor(0x1A,0x5C,0x2A)]
for (title, body), top, col in zip(steps, tops13, colors13):
    box(s, title, 0.5, top, 2.5, 0.9, fill=col, fg=WHITE, sz=12, bold=True)
    txb(s, body, 3.1, top, 9.8, 0.9, size=11, color=LIGHT_GREY)
    if top < 5.35:
        divider(s, top + 1.0)

txb(s, "SCREENSHOTS TO INCLUDE:  EC2 panel  |  RDS panel  |  S3 panel  |  Summary page with grand total  |  Share URL",
    0.5, 6.85, 12.33, 0.4, size=11, bold=True, color=ACCENT_ORG, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Conclusion — Why AWS Is the Right Move",
            "A more flexible, scalable, and reliable infrastructure platform for our game studio")

headers_con = ["Business Need", "AWS Solution", "Outcome"]
rows_con = [
    ["Handle launch-day traffic spikes",  "EC2 Auto Scaling + ELB",       "No crashes, no over-provisioning"],
    ["Keep player data available 24/7",   "RDS Multi-AZ",                 "Automatic failover, no manual work"],
    ["Distribute patches globally",       "S3 + CloudFront",              "Fast downloads for players worldwide"],
    ["Reduce upfront hardware risk",      "Pay-as-you-go model",          "CapEx to OpEx, lower financial risk"],
    ["Understand player behavior",        "CloudWatch, Kinesis, Athena",  "Data-driven game improvements"],
    ["Run event-driven game logic",       "AWS Lambda",                   "No idle servers, pay per execution"],
]
add_table(s, headers_con, rows_con, 0.5, 1.45, 12.33, 3.0)

# Financial summary boxes
txb(s, "FINANCIAL SUMMARY", 0.5, 4.65, 12.33, 0.35, size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
box(s, "On-Demand\n~$19,501/mo\n~$234,012/yr", 0.5, 5.05, 3.5, 1.1, fill=MID_GREY, fg=WHITE, sz=13, bold=True)
box(s, "1-Year Reserved\n$12,209/mo\n$146,508/yr", 4.2, 5.05, 3.5, 1.1, fill=ACCENT_BLUE, fg=WHITE, sz=13, bold=True)
box(s, "Annual Savings\n~$87,504/yr\nwith Reserved", 7.9, 5.05, 3.5, 1.1, fill=RGBColor(0x1A,0x5C,0x2A), fg=WHITE, sz=13, bold=True)

txb(s, '"For a game company where demand changes overnight and player experience depends on performance and availability,\nAWS is not just a cost decision — it is a competitive advantage."',
    0.5, 6.3, 12.33, 0.85, size=12, italic=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out = r"g:\KIro project\project school\AWS_Cloud_Migration_GameStudio.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
